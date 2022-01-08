from PyQt5 import QtCore, QtGui, QtWidgets
from google_trans_new import google_translator, LANGUAGES
import pptx
from PyQt5.QtWidgets import QMessageBox
from tkinter.filedialog import askopenfilename
from tkinter import Tk, ttk

class Ui_MainWindow(object):
    def setupUi(self, MainWindow):
        MainWindow.setObjectName("Translator")
        MainWindow.resize(380, 200)
        MainWindow.setMinimumSize(QtCore.QSize(380, 200))
        MainWindow.setMaximumSize(QtCore.QSize(380, 200))
        self.centralwidget = QtWidgets.QWidget(MainWindow)
        self.centralwidget.setObjectName("centralwidget")
        self.gridLayoutWidget = QtWidgets.QWidget(self.centralwidget)
        self.gridLayoutWidget.setGeometry(QtCore.QRect(10, 10, 361, 121))
        self.gridLayoutWidget.setObjectName("gridLayoutWidget")
        self.gridLayout = QtWidgets.QGridLayout(self.gridLayoutWidget)
        self.gridLayout.setContentsMargins(0, 0, 0, 0)
        self.gridLayout.setObjectName("gridLayout")
        self.lblName = QtWidgets.QLabel(self.gridLayoutWidget)
        self.lblName.setObjectName("lblName")
        self.gridLayout.addWidget(self.lblName, 4, 1, 1, 1)
        self.label_3 = QtWidgets.QLabel(self.gridLayoutWidget)
        self.label_3.setObjectName("label_3")
        self.gridLayout.addWidget(self.label_3, 4, 0, 1, 1)
        self.listLanguages = QtWidgets.QComboBox(self.gridLayoutWidget)
        self.listLanguages.setMinimumSize(QtCore.QSize(250, 0))
        self.listLanguages.setObjectName("listLanguages")
        self.gridLayout.addWidget(self.listLanguages, 2, 1, 1, 1)
        self.label = QtWidgets.QLabel(self.gridLayoutWidget)
        self.label.setObjectName("label")
        self.gridLayout.addWidget(self.label, 2, 0, 1, 1)
        self.label_2 = QtWidgets.QLabel(self.gridLayoutWidget)
        self.label_2.setMinimumSize(QtCore.QSize(0, 50))
        self.label_2.setMaximumSize(QtCore.QSize(16777215, 50))
        self.label_2.setObjectName("label_2")
        self.gridLayout.addWidget(self.label_2, 0, 0, 1, 2)
        self.horizontalLayoutWidget = QtWidgets.QWidget(self.centralwidget)
        self.horizontalLayoutWidget.setGeometry(QtCore.QRect(10, 130, 361, 61))
        self.horizontalLayoutWidget.setObjectName("horizontalLayoutWidget")
        self.horizontalLayout = QtWidgets.QHBoxLayout(self.horizontalLayoutWidget)
        self.horizontalLayout.setContentsMargins(0, 0, 0, 0)
        self.horizontalLayout.setObjectName("horizontalLayout")
        self.btnBrowse = QtWidgets.QPushButton(self.horizontalLayoutWidget)
        self.btnBrowse.setMinimumSize(QtCore.QSize(0, 50))
        icon = QtGui.QIcon()
        icon.addPixmap(QtGui.QPixmap("Resources/search.png"), QtGui.QIcon.Normal, QtGui.QIcon.Off)
        self.btnBrowse.setIcon(icon)
        self.btnBrowse.setObjectName("btnBrowse")
        self.horizontalLayout.addWidget(self.btnBrowse)
        self.btnTranslate = QtWidgets.QPushButton(self.horizontalLayoutWidget)
        self.btnTranslate.setMinimumSize(QtCore.QSize(0, 50))
        icon1 = QtGui.QIcon()
        icon1.addPixmap(QtGui.QPixmap("Resources/translate.png"), QtGui.QIcon.Normal, QtGui.QIcon.Off)
        self.btnTranslate.setIcon(icon1)
        self.btnTranslate.setObjectName("btnTranslate")
        self.horizontalLayout.addWidget(self.btnTranslate)
        MainWindow.setCentralWidget(self.centralwidget)
        self.retranslateUi(MainWindow)
        QtCore.QMetaObject.connectSlotsByName(MainWindow)
        self.btnBrowse.clicked.connect(lambda: getFile(self.lblName))
        self.btnTranslate.clicked.connect(lambda: translate(self.listLanguages, self.lblName))

    def retranslateUi(self, MainWindow):
        _translate = QtCore.QCoreApplication.translate
        MainWindow.setWindowTitle(_translate("MainWindow", "PowerPoint Tanslator"))
        self.lblName.setText(_translate("MainWindow", "Unloaded"))
        self.label_3.setText(_translate("MainWindow", "FILE NAME:"))
        self.label.setText(_translate("MainWindow", "LANGUAGE:"))
        self.label_2.setText(_translate("MainWindow", "<html><head/><body><p align=\"center\"><span style=\" font-size:18pt; font-weight:600; color:#3465a4;\">PRESENTATION TRANSLATOR</span></p></body></html>"))
        self.btnBrowse.setText(_translate("MainWindow", " BROWSE"))
        self.btnTranslate.setText(_translate("MainWindow", " TRANSLATE"))

        for lang in LANGUAGES:
            key = (f'{lang}')
            language = (f'{LANGUAGES[lang]}')
            self.listLanguages.addItem(language)

def getFile(name):
    Tk().withdraw()
    fname = askopenfilename(initialdir="/",filetypes =(("Powerpoint", "*.pptx"),("All Files","*.*")),title = "Choose a file.")
    fname = str(fname)
    fname = fname.replace('(','')
    fname = fname.replace(')', '')
    fname = fname.replace("'", '')
    fname = fname.replace(',', '')
    name.setText(fname)

def translate(selectedLanguage, name):
    window = Tk()
    window.title("Translating")
    bar = ttk.Progressbar(window, length = 300, mode = 'indeterminate')
    bar.pack(pady=20, padx=20)

    if name.text() == "Unloaded":
        showmsg()
        window.destroy()
        return
    elif name.text() == "":
        showmsg()
        window.destroy()
        return
    else:
        translator = google_translator()
        prs = pptx.Presentation(name.text())

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for p in shape.text_frame.paragraphs:
                    formats, newRuns = [], []

                    for _, r in enumerate(p.runs):
                        text = r.text
                        getLang = selectedLanguage.currentText()
                        key = ""

                        for lang in LANGUAGES:
                            if getLang == (f'{LANGUAGES[lang]}'):
                                key = (f'{lang}')

                        bar['value'] += 20
                        window.update()
                        translated = translator.translate(text, lang_tgt=key)
                        newRuns.append(translated)
                        formats.append({'size': r.font.size,
                                        'bold': r.font.bold,
                                        'underline': r.font.underline,
                                        'italic': r.font.italic})
                    p.clear()
                    for i in range(len(newRuns)):
                        run = p.add_run()
                        run.text = newRuns[i]
                        run.font.bold = formats[i]['bold']
                        run.font.italic = formats[i]['italic']
                        run.font.size = formats[i]['size']
                        run.font.underline = formats[i]['underline']

        window.destroy()
        name = str(name.text())
        name = name.replace('.pptx', '')
        prs.save(name + ' translated to ' + selectedLanguage.currentText() + '.pptx')
        completeMsg()

def showmsg():
    msg = QMessageBox()
    msg.setIcon(QMessageBox.Warning)
    msg.setText("Please Browse for a presentation you would like to translate")
    msg.setWindowTitle("Error")
    msg.setStandardButtons(QMessageBox.Ok)
    retval = msg.exec_()

def completeMsg():
    msg = QMessageBox()
    msg.setIcon(QMessageBox.Information)
    msg.setText("Translation Complete. Please see the new file save in the location of the existing file")
    msg.setWindowTitle("Complete")
    msg.setStandardButtons(QMessageBox.Ok)
    retval = msg.exec_()

if __name__ == "__main__":
    import sys
    app = QtWidgets.QApplication(sys.argv)
    MainWindow = QtWidgets.QMainWindow()
    ui = Ui_MainWindow()
    ui.setupUi(MainWindow)
    MainWindow.show()
    sys.exit(app.exec_())